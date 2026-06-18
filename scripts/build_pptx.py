"""
Generate team-presentation-v4.5.3.pptx — visual redesign.
Replaces text-heavy bullets with flowcharts, stat cards, timeline,
and minimal copy. Instrumental theme: green left bar + footer.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------- Palette ----------
INST_GREEN   = RGBColor(0x19, 0xA9, 0x74)
INST_GREEN_DK= RGBColor(0x05, 0x96, 0x69)
INST_DARK    = RGBColor(0x0F, 0x17, 0x2A)
INST_MUTED   = RGBColor(0x64, 0x74, 0x8B)
INST_LIGHT   = RGBColor(0x94, 0xA3, 0xB8)
INST_BG      = RGBColor(0xF8, 0xFA, 0xFC)
INST_BG_2    = RGBColor(0xF1, 0xF5, 0xF9)
INST_BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLANK_BG     = RGBColor(0xFE, 0xF3, 0xC7)
BLANK_FG     = RGBColor(0x92, 0x40, 0x0E)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
BLUE_BG      = RGBColor(0xEF, 0xF6, 0xFF)
GREEN_LIGHT  = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT_BG = RGBColor(0xEC, 0xFD, 0xF5)
RED          = RGBColor(0xDC, 0x26, 0x26)
RED_BG       = RGBColor(0xFE, 0xF2, 0xF2)
ORANGE       = RGBColor(0xD9, 0x77, 0x06)
ORANGE_BG    = RGBColor(0xFF, 0xFB, 0xEB)
PINK         = RGBColor(0xDB, 0x27, 0x77)
PINK_BG      = RGBColor(0xFC, 0xE7, 0xF3)
PURPLE       = RGBColor(0xA8, 0x55, 0xF7)
PURPLE_BG    = RGBColor(0xFA, 0xF5, 0xFF)
TEAL         = RGBColor(0x0D, 0x94, 0x88)
TEAL_BG      = RGBColor(0xCC, 0xFB, 0xF1)

# ---------- Layout ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BAR_W   = Emu(120000)
LEFT_M  = Inches(0.6)
RIGHT_M = Inches(0.55)
TOP_M   = Inches(0.45)
BOT_M   = Inches(0.5)

# ---------- Helpers ----------
def _no_shadow(shape):
    shape.shadow.inherit = False

def add_bar_and_footer(slide, page_num):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, BAR_W, SLIDE_H)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = INST_GREEN
    _no_shadow(bar)
    # divider line near footer
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, LEFT_M, SLIDE_H - BOT_M - Inches(0.22),
                                       SLIDE_W - RIGHT_M, SLIDE_H - BOT_M - Inches(0.22))
    line.line.color.rgb = INST_BORDER
    line.line.width = Pt(0.5)
    # center text
    tb = slide.shapes.add_textbox(LEFT_M, SLIDE_H - BOT_M - Inches(0.18), SLIDE_W - LEFT_M - RIGHT_M, Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Confidential — Instrumental"
    r.font.size = Pt(9); r.font.color.rgb = INST_LIGHT
    # logo right
    tb2 = slide.shapes.add_textbox(SLIDE_W - RIGHT_M - Inches(2.2), SLIDE_H - BOT_M - Inches(0.18), Inches(1.7), Inches(0.28))
    tf2 = tb2.text_frame
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = "INSTRUMENTAL"
    r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = INST_DARK
    # page num
    tb3 = slide.shapes.add_textbox(SLIDE_W - RIGHT_M - Inches(0.5), SLIDE_H - BOT_M - Inches(0.18), Inches(0.45), Inches(0.28))
    tf3 = tb3.text_frame
    tf3.margin_left = tf3.margin_right = tf3.margin_top = tf3.margin_bottom = 0
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.RIGHT
    r3 = p3.add_run(); r3.text = str(page_num)
    r3.font.size = Pt(10); r3.font.bold = True; r3.font.color.rgb = INST_MUTED

def add_title(slide, title, subtitle=None, size=32):
    tb = slide.shapes.add_textbox(LEFT_M, TOP_M, SLIDE_W - LEFT_M - RIGHT_M, Inches(0.6))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(size); r.font.color.rgb = INST_DARK
    if subtitle:
        sub = slide.shapes.add_textbox(LEFT_M, TOP_M + Inches(0.55), SLIDE_W - LEFT_M - RIGHT_M, Inches(0.35))
        tf = sub.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = subtitle
        r.font.size = Pt(14); r.font.color.rgb = INST_MUTED

def add_text(slide, text, x, y, w, h, size=11, color=INST_DARK, bold=False, align=PP_ALIGN.LEFT,
              italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    return tb

def add_card(slide, x, y, w, h, fill, border, border_w=1.5, radius=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    _no_shadow(sh)
    return sh

def add_pill(slide, label, x, y, w, h, fill, border, fg=None, size=10, bold=True):
    sh = add_card(slide, x, y, w, h, fill, border, border_w=1.2)
    tf = sh.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = fg if fg else INST_DARK
    return sh

def add_arrow(slide, x1, y1, x2, y2, color=INST_LIGHT, weight=1.5):
    arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    arrow.line.color.rgb = color
    arrow.line.width = Pt(weight)
    # set arrow head
    ln = arrow.line._get_or_add_ln()
    tail_end = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail_end)
    return arrow

def add_arrow_horiz(slide, x1, y1, x2, y2, color=INST_GREEN, weight=2.0):
    """Right-pointing horizontal arrow."""
    return add_arrow(slide, x1, y1, x2, y2, color=color, weight=weight)

def card_with_title_body(slide, x, y, w, h, title, body, fill=INST_BG, border=INST_BORDER,
                          title_color=INST_DARK, title_size=14, body_size=10.5, accent=None):
    sh = add_card(slide, x, y, w, h, fill, border)
    if accent:
        # accent stripe on the left of the card
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        _no_shadow(stripe)
    tf = sh.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.10)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(title_size); r.font.bold = True; r.font.color.rgb = title_color
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    _render_text_with_blanks(p2, body, size=body_size)

def _render_text_with_blanks(paragraph, text, size=10.5, color=INST_DARK):
    """Render text where [BLANK ...] is highlighted yellow/italic."""
    import re
    parts = re.split(r'(\[BLANK[^\]]*\])', text)
    for part in parts:
        if not part: continue
        r = paragraph.add_run()
        r.text = part
        r.font.size = Pt(size)
        if part.startswith("[BLANK"):
            r.font.color.rgb = BLANK_FG
            r.font.italic = True
        else:
            r.font.color.rgb = color

def stat_card(slide, x, y, w, h, big_text, small_text, fill=INST_BG, accent=INST_GREEN):
    sh = add_card(slide, x, y, w, h, fill, accent, border_w=2)
    tf = sh.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big_text
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = accent
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(0)
    r2 = p2.add_run(); r2.text = small_text
    r2.font.size = Pt(9.5); r2.font.color.rgb = INST_DARK

def labeled_box(slide, x, y, w, h, label, sub=None, fill=WHITE, border=INST_DARK, size=11):
    sh = add_card(slide, x, y, w, h, fill, border, border_w=1.5)
    tf = sh.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INST_DARK
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(8.5); r2.font.color.rgb = INST_MUTED
    return sh

# ====================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# =========== SLIDE 1 — HERO ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 1)
add_title(s, "Deployment Portal",
          "A PMO frontend that turns HubSpot into a deployment-management tool", size=40)

# Hero flow: User -> App -> HubSpot
flow_y = TOP_M + Inches(1.4)
flow_h = Inches(1.0)
box_w = Inches(2.6)
gap = Inches(0.7)
start_x = LEFT_M + Inches(0.8)

labeled_box(s, start_x, flow_y, box_w, flow_h, "CX Team",
            "Instrumental + SI", fill=INST_BG_2, border=INST_MUTED, size=15)
add_arrow(s, start_x + box_w + Inches(0.05), flow_y + flow_h/2,
          start_x + box_w + gap - Inches(0.05), flow_y + flow_h/2,
          color=INST_GREEN, weight=2.5)
# webapp box (highlight)
labeled_box(s, start_x + box_w + gap, flow_y, box_w, flow_h, "Deployment Portal",
            "React + Firebase", fill=GREEN_LIGHT_BG, border=INST_GREEN, size=15)
# arrows out — bidirectional
add_arrow(s, start_x + box_w*2 + gap + Inches(0.05), flow_y + flow_h/2,
          start_x + box_w*2 + gap*2 - Inches(0.05), flow_y + flow_h/2,
          color=INST_GREEN, weight=2.5)
add_arrow(s, start_x + box_w*2 + gap*2 - Inches(0.05), flow_y + flow_h/2 + Inches(0.15),
          start_x + box_w*2 + gap + Inches(0.05), flow_y + flow_h/2 + Inches(0.15),
          color=INST_GREEN, weight=2.5)
# back-arrow caption
add_text(s, "bidirectional sync", start_x + box_w + gap + Inches(0.1), flow_y + flow_h + Inches(0.05),
         box_w + gap, Inches(0.25), size=9, color=INST_MUTED, italic=True, align=PP_ALIGN.CENTER)

labeled_box(s, start_x + box_w*2 + gap*2, flow_y, box_w, flow_h, "HubSpot + APIs",
            "the system of record", fill=INST_BG_2, border=INST_MUTED, size=15)

# 3 cards: What | Why | How (Two workflows)
card_y = flow_y + flow_h + Inches(0.7)
card_h = Inches(2.1)
card_w = (SLIDE_W - LEFT_M - RIGHT_M - Inches(0.5)) / 3
gap_c = Inches(0.25)

card_with_title_body(s, LEFT_M + Inches(0.1), card_y, card_w, card_h,
    "What", "Consolidated UI for projects, station kits, fleet assets, shipments, DRIs — bidirectional with HubSpot.",
    fill=WHITE, border=INST_BORDER, accent=INST_GREEN, title_size=18, body_size=12)
card_with_title_body(s, LEFT_M + Inches(0.1) + card_w + gap_c, card_y, card_w, card_h,
    "Why", "HubSpot alone = hard to share externally, not AI-native, clunky for PMO work, no home for non-CRM data.",
    fill=WHITE, border=INST_BORDER, accent=BLUE, title_size=18, body_size=12)
card_with_title_body(s, LEFT_M + Inches(0.1) + (card_w + gap_c)*2, card_y, card_w, card_h,
    "Two workflows", "Instrumental Hardware Deployment (Asang) + SI Partner Deployment (Sneha). Same app today; splitting next.",
    fill=WHITE, border=INST_BORDER, accent=PURPLE, title_size=18, body_size=12)

# API pills row
pill_y = card_y + card_h + Inches(0.25)
api_labels = ["HubSpot", "Firebase RTDB", "Firebase Auth", "Cloud Functions", "Claude API", "Slack"]
pill_w = Inches(1.55); pill_h = Inches(0.32); gap_p = Inches(0.1)
total_pills_w = pill_w * len(api_labels) + gap_p * (len(api_labels)-1)
start_pills = LEFT_M + (SLIDE_W - LEFT_M - RIGHT_M - total_pills_w) / 2
for i, lbl in enumerate(api_labels):
    add_pill(s, lbl, start_pills + (pill_w + gap_p)*i, pill_y, pill_w, pill_h,
             fill=INST_BG, border=INST_BORDER, size=10)

# =========== SLIDE 2 — TOC (visual grid) ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 2)
add_title(s, "What we'll cover")

toc_items = [
    ("What & Why", "the webapp + reason built"),
    ("Vision", "PLM, AI OS, workflow split"),
    ("Hosting & Auth", "free vs paid + JWT"),
    ("Under the hood", "frontend + backend"),
    ("Build & ship", "Claude Code + workflow + history"),
    ("Ops integration", "where it fits + what to connect"),
]
grid_y = TOP_M + Inches(1.3)
cols = 3; rows = 2
gw = (SLIDE_W - LEFT_M - RIGHT_M - Inches(0.4))/cols - Inches(0.2)
gh = Inches(2.1)
gap_x = Inches(0.3); gap_y = Inches(0.3)
accent_palette = [INST_GREEN, BLUE, PURPLE, PINK, ORANGE, TEAL]
for i, (lab, sub) in enumerate(toc_items):
    r = i // cols; c = i % cols
    x = LEFT_M + Inches(0.2) + c * (gw + gap_x)
    y = grid_y + r * (gh + gap_y)
    sh = add_card(s, x, y, gw, gh, WHITE, INST_BORDER, border_w=1.5)
    # number badge
    num = add_card(s, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.55),
                   accent_palette[i], accent_palette[i], border_w=0)
    tf = num.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = str(i+1)
    rr.font.size = Pt(22); rr.font.bold = True; rr.font.color.rgb = WHITE
    # title + sub
    add_text(s, lab, x + Inches(0.95), y + Inches(0.25), gw - Inches(1.15), Inches(0.45),
             size=18, bold=True, color=INST_DARK)
    add_text(s, sub, x + Inches(0.95), y + Inches(0.75), gw - Inches(1.15), Inches(1.0),
             size=11.5, color=INST_MUTED)

# =========== SLIDE 3 — Vision (3-card timeline) ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 3)
add_title(s, "Where this goes next", size=32)

# horizontal timeline with 3 cards: Today | Near-term | Long-term
tl_y = TOP_M + Inches(1.4)
tl_h = Inches(3.2)
tl_w = (SLIDE_W - LEFT_M - RIGHT_M - Inches(0.6)) / 3
gap_tl = Inches(0.3)
start_tl = LEFT_M + Inches(0.2)

# Today
card_with_title_body(s, start_tl, tl_y, tl_w, tl_h,
    "Today",
    "Single app, two workflows muddled (Instrumental + SI). Heavy HubSpot bidirectional sync. v4.5.x in production.",
    fill=INST_BG, border=INST_BORDER, accent=INST_MUTED, title_size=20, body_size=12)
# Near-term
card_with_title_body(s, start_tl + tl_w + gap_tl, tl_y, tl_w, tl_h,
    "Near-term",
    "Split UIs (same backend) — Instrumental + SI each get optimized views. Adoption push across CX. Connect adjacent workstreams.",
    fill=GREEN_LIGHT_BG, border=INST_GREEN, accent=INST_GREEN, title_size=20, body_size=12)
# Long-term
card_with_title_body(s, start_tl + (tl_w + gap_tl)*2, tl_y, tl_w, tl_h,
    "Long-term",
    "Full PLM tool — replace Odoo / Arena for inventory, BOM, ECO. Integrate with Instrumental AI OS.\n\n[BLANK — AI OS specifics]",
    fill=PURPLE_BG, border=PURPLE, accent=PURPLE, title_size=20, body_size=12)

# Arrows between cards
arrow_y = tl_y + tl_h/2
add_arrow(s, start_tl + tl_w + Inches(0.05), arrow_y,
          start_tl + tl_w + gap_tl - Inches(0.05), arrow_y, color=INST_GREEN, weight=3)
add_arrow(s, start_tl + (tl_w + gap_tl)*2 - gap_tl + Inches(0.05), arrow_y,
          start_tl + (tl_w + gap_tl)*2 - Inches(0.05), arrow_y, color=PURPLE, weight=3)

# Footer note re: workflow split rationale
note_y = tl_y + tl_h + Inches(0.3)
add_card(s, LEFT_M + Inches(0.2), note_y, SLIDE_W - LEFT_M - RIGHT_M - Inches(0.4), Inches(0.7),
        INST_BG, INST_BORDER)
add_text(s, "Why split the UIs?  ", LEFT_M + Inches(0.4), note_y + Inches(0.18), Inches(2), Inches(0.3),
         size=12, bold=True, color=INST_DARK)
add_text(s, "Different stages, milestone definitions, and stakeholders → one merged UI adds confusion. Sneha owns SI; Asang owns Instrumental; backend stays shared.",
         LEFT_M + Inches(2.1), note_y + Inches(0.15), SLIDE_W - LEFT_M - RIGHT_M - Inches(2.7), Inches(0.5),
         size=11, color=INST_DARK)

# =========== SLIDE 4 — Hosting + JWT ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 4)
add_title(s, "Hosting + auth — free today, when we'd upgrade")

# LEFT: Stack diagram (vertical)
stack_x = LEFT_M + Inches(0.2)
stack_y = TOP_M + Inches(1.2)
stack_w = Inches(3.6)
layer_h = Inches(0.55)
layer_gap = Inches(0.12)
layers = [
    ("Browser (React 18 + Vite)", WHITE, INST_DARK),
    ("Firebase Hosting", INST_BG, INST_MUTED),
    ("Firebase Auth (Google OAuth)", GREEN_LIGHT_BG, INST_GREEN),
    ("Cloud Functions (Node 20)", BLUE_BG, BLUE),
    ("Realtime Database (RTDB)", ORANGE_BG, ORANGE),
]
for i, (lbl, fill, border) in enumerate(layers):
    y = stack_y + i * (layer_h + layer_gap)
    labeled_box(s, stack_x, y, stack_w, layer_h, lbl, fill=fill, border=border, size=12)
# arrow on the side: external APIs (HubSpot, Claude, Slack)
ext_x = stack_x + stack_w + Inches(0.4)
ext_y = stack_y + 3*(layer_h+layer_gap)
labeled_box(s, ext_x, ext_y, Inches(1.9), layer_h, "HubSpot · Claude · Slack",
            fill=INST_BG, border=INST_MUTED, size=10)
add_arrow(s, stack_x + stack_w + Inches(0.05), ext_y + layer_h/2,
          ext_x - Inches(0.05), ext_y + layer_h/2, color=BLUE, weight=2)

# Stack title
add_text(s, "Stack", stack_x, stack_y - Inches(0.4), stack_w, Inches(0.3),
         size=13, bold=True, color=INST_DARK)

# Cost badges below stack
cost_y = stack_y + len(layers)*(layer_h + layer_gap) + Inches(0.25)
badges = [
    ("Spark (free) today", GREEN_LIGHT_BG, INST_GREEN),
    ("Blaze if > 2M CF/mo", ORANGE_BG, ORANGE),
    ("AWS — [BLANK]", BLANK_BG, BLANK_FG),
]
bw = Inches(1.7); gap_b = Inches(0.15)
for i, (lbl, fill, border) in enumerate(badges):
    add_pill(s, lbl, stack_x + i*(bw + gap_b), cost_y, bw, Inches(0.4),
             fill=fill, border=border, fg=border, size=10)

# RIGHT: JWT flow (4-step horizontal)
jwt_x = LEFT_M + Inches(6.4)
jwt_w = SLIDE_W - RIGHT_M - jwt_x - Inches(0.1)
add_text(s, "JWT auth flow", jwt_x, TOP_M + Inches(0.8), jwt_w, Inches(0.3),
         size=13, bold=True, color=INST_DARK)

step_y = TOP_M + Inches(1.2)
step_h = Inches(0.85)
step_gap = Inches(0.18)
steps = [
    ("1", "Sign in with Google", INST_GREEN),
    ("2", "Firebase Auth issues JWT", BLUE),
    ("3", "JWT on every request", PURPLE),
    ("4", "Rules verify role at DB layer", INST_GREEN_DK),
]
for i, (num, txt, color) in enumerate(steps):
    y = step_y + i*(step_h + step_gap)
    sh = add_card(s, jwt_x, y, jwt_w, step_h, WHITE, color, border_w=2)
    # number badge
    nb = add_card(s, jwt_x + Inches(0.15), y + Inches(0.15), Inches(0.55), step_h - Inches(0.3),
                  color, color, border_w=0)
    tf = nb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = num
    rr.font.size = Pt(20); rr.font.bold = True; rr.font.color.rgb = WHITE
    add_text(s, txt, jwt_x + Inches(0.85), y, jwt_w - Inches(1.0), step_h,
             size=13, bold=True, color=INST_DARK, anchor=MSO_ANCHOR.MIDDLE)

# =========== SLIDE 5 — Tech architecture ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 5)
add_title(s, "Under the hood")

# Architecture flowchart (top half)
arch_y = TOP_M + Inches(1.0)
arch_h = Inches(3.3)

# Browser (left)
labeled_box(s, LEFT_M + Inches(0.3), arch_y + Inches(0.8), Inches(2.1), Inches(0.9),
            "Browser", "React 18 + Vite\n~770 KB bundle",
            fill=INST_BG_2, border=INST_DARK, size=13)
# Arrow → Hosting
arrow_y = arch_y + Inches(1.25)
add_arrow(s, LEFT_M + Inches(2.45), arrow_y, LEFT_M + Inches(2.95), arrow_y,
          color=INST_GREEN, weight=2.5)
# Firebase Hosting
labeled_box(s, LEFT_M + Inches(3.0), arch_y + Inches(0.5), Inches(2.4), Inches(0.6),
            "Firebase Hosting", fill=GREEN_LIGHT_BG, border=INST_GREEN, size=12)
# Cloud Functions
labeled_box(s, LEFT_M + Inches(3.0), arch_y + Inches(1.4), Inches(2.4), Inches(0.6),
            "Cloud Functions · Node 20", "20+ callables", fill=BLUE_BG, border=BLUE, size=12)
# Auth
labeled_box(s, LEFT_M + Inches(3.0), arch_y + Inches(2.3), Inches(2.4), Inches(0.6),
            "Firebase Auth", "Google OAuth + JWT", fill=PURPLE_BG, border=PURPLE, size=12)
# arrows from browser to each
add_arrow(s, LEFT_M + Inches(2.45), arrow_y, LEFT_M + Inches(2.95), arch_y + Inches(0.8), color=INST_GREEN, weight=2)
add_arrow(s, LEFT_M + Inches(2.45), arrow_y, LEFT_M + Inches(2.95), arch_y + Inches(1.7), color=BLUE, weight=2)
add_arrow(s, LEFT_M + Inches(2.45), arrow_y, LEFT_M + Inches(2.95), arch_y + Inches(2.6), color=PURPLE, weight=2)

# RTDB
labeled_box(s, LEFT_M + Inches(6.0), arch_y + Inches(0.95), Inches(2.5), Inches(0.9),
            "Realtime Database", "projects · docData · users\nsync log · audit log",
            fill=ORANGE_BG, border=ORANGE, size=13)
add_arrow(s, LEFT_M + Inches(5.45), arch_y + Inches(1.7), LEFT_M + Inches(5.95), arch_y + Inches(1.4), color=ORANGE, weight=2)
add_arrow(s, LEFT_M + Inches(5.95), arch_y + Inches(1.4), LEFT_M + Inches(5.45), arch_y + Inches(1.7), color=ORANGE, weight=2)

# External APIs
labeled_box(s, LEFT_M + Inches(9.2), arch_y + Inches(0.5), Inches(2.8), Inches(0.6),
            "HubSpot APIs", "Custom Objects · Owners · v4 Assoc.",
            fill=PINK_BG, border=PINK, size=11)
labeled_box(s, LEFT_M + Inches(9.2), arch_y + Inches(1.4), Inches(2.8), Inches(0.6),
            "Anthropic Claude API", "in-app AI bot",
            fill=GREEN_LIGHT_BG, border=GREEN_LIGHT, size=11)
labeled_box(s, LEFT_M + Inches(9.2), arch_y + Inches(2.3), Inches(2.8), Inches(0.6),
            "Slack webhooks", "feedback routing",
            fill=TEAL_BG, border=TEAL, size=11)
# arrows from CFs to external
add_arrow(s, LEFT_M + Inches(5.45), arch_y + Inches(1.7), LEFT_M + Inches(9.15), arch_y + Inches(0.8), color=PINK, weight=1.8)
add_arrow(s, LEFT_M + Inches(5.45), arch_y + Inches(1.7), LEFT_M + Inches(9.15), arch_y + Inches(1.7), color=GREEN_LIGHT, weight=1.8)
add_arrow(s, LEFT_M + Inches(5.45), arch_y + Inches(1.7), LEFT_M + Inches(9.15), arch_y + Inches(2.6), color=TEAL, weight=1.8)

# Stat callouts at bottom
stat_y = arch_y + arch_h + Inches(0.25)
stat_h = Inches(1.2)
stats = [
    ("20+", "Cloud Functions"),
    ("~11K", "lines in App.jsx"),
    ("5", "languages supported"),
    ("DB-level", "access control"),
]
sw = (SLIDE_W - LEFT_M - RIGHT_M - Inches(0.6)) / 4
gap_s = Inches(0.2)
for i, (big, small) in enumerate(stats):
    x = LEFT_M + Inches(0.3) + i*(sw + gap_s)
    stat_card(s, x, stat_y, sw, stat_h, big, small,
              fill=INST_BG, accent=[INST_GREEN, BLUE, PURPLE, ORANGE][i])

# =========== SLIDE 6 — Build & ship ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 6)
add_title(s, "Built with Claude Code — workflow + history")

# TOP: 5-step release workflow (horizontal)
flow_y = TOP_M + Inches(1.05)
step_w = (SLIDE_W - LEFT_M - RIGHT_M - Inches(0.5)) / 5 - Inches(0.1)
gap_f = Inches(0.1)
flow_steps = [
    ("Bump", "version + README"),
    ("Build", "npm run build"),
    ("Deploy", "firebase deploy"),
    ("Verify", "manual sync + RTDB"),
    ("Ship", "commit + zip + push"),
]
for i, (lbl, sub) in enumerate(flow_steps):
    x = LEFT_M + Inches(0.3) + i*(step_w + gap_f)
    sh = add_card(s, x, flow_y, step_w, Inches(1.0), WHITE, INST_GREEN, border_w=2)
    add_text(s, lbl, x, flow_y + Inches(0.18), step_w, Inches(0.3),
             size=15, bold=True, color=INST_GREEN, align=PP_ALIGN.CENTER)
    add_text(s, sub, x, flow_y + Inches(0.5), step_w, Inches(0.4),
             size=10, color=INST_MUTED, align=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_arrow(s, x + step_w + Inches(0.005), flow_y + Inches(0.5),
                  x + step_w + gap_f - Inches(0.005), flow_y + Inches(0.5),
                  color=INST_GREEN, weight=2)

# LEFT bottom: Lessons (4 mini cards)
lessons_y = flow_y + Inches(1.3)
add_text(s, "Lessons from building with Claude Code",
         LEFT_M + Inches(0.3), lessons_y, Inches(6), Inches(0.3),
         size=13, bold=True, color=INST_DARK)
lessons = [
    ("Define components at module scope", "Nested components remount + lose state", RED),
    ("Field-level DB writes only", "Whole-node sets can wipe RTDB (v4.5.1)", ORANGE),
    ("Hardcode known type IDs", "Discovery breaks on HubSpot reorder (v4.5.2)", BLUE),
    ("Handle HubSpot 429/5xx explicitly", "Silent partial sync otherwise", PURPLE),
]
lc_y = lessons_y + Inches(0.4)
lc_w = (Inches(6.3) - Inches(0.15)) / 2
lc_h = Inches(1.0)
for i, (title, body, color) in enumerate(lessons):
    r = i // 2; c = i % 2
    x = LEFT_M + Inches(0.3) + c*(lc_w + Inches(0.15))
    y = lc_y + r*(lc_h + Inches(0.15))
    card_with_title_body(s, x, y, lc_w, lc_h, title, body,
                         fill=WHITE, border=INST_BORDER, accent=color,
                         title_size=11, body_size=9.5)

# RIGHT bottom: Version timeline
ver_x = LEFT_M + Inches(7.0)
ver_w = SLIDE_W - RIGHT_M - ver_x - Inches(0.1)
add_text(s, "Version history", ver_x, lessons_y, ver_w, Inches(0.3),
         size=13, bold=True, color=INST_DARK)
versions = [
    ("v3.0", "Security review PASSED"),
    ("v3.x", "HubSpot sync · DB-level access"),
    ("v4.0", "Security response · Project Overview"),
    ("v4.1", "Date writeback · si_admin"),
    ("v4.4", "Shipment writeback"),
    ("v4.5", "Kit + Fleet Asset writeback · DRI sync"),
]
v_y = lessons_y + Inches(0.4)
v_h = Inches(0.36)
v_gap = Inches(0.04)
for i, (ver, desc) in enumerate(versions):
    y = v_y + i*(v_h + v_gap)
    # dot
    dot = add_card(s, ver_x, y + Inches(0.06), Inches(0.22), Inches(0.22),
                   INST_GREEN, INST_GREEN, border_w=0)
    # ver
    add_text(s, ver, ver_x + Inches(0.32), y + Inches(0.03), Inches(0.6), Inches(0.3),
             size=11, bold=True, color=INST_DARK)
    # desc
    add_text(s, desc, ver_x + Inches(0.95), y + Inches(0.03), ver_w - Inches(1.0), Inches(0.3),
             size=10, color=INST_MUTED)

# Security strip
sec_y = lc_y + 2*(lc_h + Inches(0.15)) + Inches(0.1)
sec = add_card(s, LEFT_M + Inches(0.3), sec_y, Inches(6.3), Inches(0.5),
               GREEN_LIGHT_BG, INST_GREEN, border_w=1.5)
add_text(s, "🔒  Security review v3.0.0 passed · DB-level access · audit log · backups enabled",
         LEFT_M + Inches(0.45), sec_y, Inches(6.1), Inches(0.5),
         size=10.5, color=INST_GREEN_DK, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# =========== SLIDE 7 — OPS workflow ===========
s = prs.slides.add_slide(blank)
add_bar_and_footer(s, 7)
add_title(s, "Where this fits + what's next to connect")

# LEFT: diagram (larger)
d_left = LEFT_M + Inches(0.2)
d_top = TOP_M + Inches(0.95)
d_w = Inches(7.2)
d_h = Inches(4.7)

# Annotation box (dashed green)
annot = add_card(s, d_left + Inches(2.0), d_top + Inches(0.25), Inches(5.0), Inches(3.05),
                 RGBColor(0xE8, 0xF7, 0xEF), INST_GREEN, border_w=2.5)
annot.line.dash_style = 7
# Label
add_text(s, "Deployment Portal (this webapp) — reads + writes",
         d_left + Inches(2.0), d_top + Inches(0.0), Inches(5.0), Inches(0.25),
         size=11.5, bold=True, color=INST_GREEN_DK, align=PP_ALIGN.CENTER)

# Object boxes inside the annotation
labeled_box(s, d_left + Inches(2.2), d_top + Inches(0.5), Inches(4.6), Inches(0.5),
            "Projects (5 types)", fill=RED_BG, border=RED, size=12)
labeled_box(s, d_left + Inches(2.2), d_top + Inches(1.15), Inches(2.2), Inches(0.7),
            "Station Kits", "kit SN · computer SN · status",
            fill=ORANGE_BG, border=ORANGE, size=12)
labeled_box(s, d_left + Inches(4.6), d_top + Inches(1.15), Inches(2.2), Inches(0.7),
            "Fleet Assets", "Camera · Lens · Computer · LED SNs",
            fill=ORANGE_BG, border=ORANGE, size=12)
labeled_box(s, d_left + Inches(2.2), d_top + Inches(2.05), Inches(4.6), Inches(0.7),
            "Shipments", "INxxx · carrier · tracking · ship date",
            fill=PINK_BG, border=PINK, size=12)
# Arrows
add_arrow(s, d_left + Inches(4.5), d_top + Inches(1.0), d_left + Inches(3.3), d_top + Inches(1.15), color=INST_LIGHT, weight=1.5)
add_arrow(s, d_left + Inches(4.5), d_top + Inches(1.0), d_left + Inches(5.7), d_top + Inches(1.15), color=INST_LIGHT, weight=1.5)
add_arrow(s, d_left + Inches(3.3), d_top + Inches(1.85), d_left + Inches(4.5), d_top + Inches(2.05), color=INST_LIGHT, weight=1.5)
add_arrow(s, d_left + Inches(5.7), d_top + Inches(1.85), d_left + Inches(4.5), d_top + Inches(2.05), color=INST_LIGHT, weight=1.5)

# OUTSIDE: SalesHub Deals (top-left)
labeled_box(s, d_left + Inches(0.0), d_top + Inches(0.5), Inches(1.8), Inches(0.55),
            "SalesHub Deals", "(no integration)", fill=GREEN_LIGHT_BG, border=GREEN_LIGHT, size=11)
# CS Programs
labeled_box(s, d_left + Inches(0.0), d_top + Inches(1.3), Inches(1.8), Inches(1.4),
            "CS Programs", "Kickoff → Onboarding\n→ Validation → Complete\n(partial integration)",
            fill=BLUE_BG, border=BLUE, size=11)
add_arrow(s, d_left + Inches(1.85), d_top + Inches(0.75), d_left + Inches(2.15), d_top + Inches(0.75), color=INST_LIGHT, weight=1.5)
add_arrow(s, d_left + Inches(1.85), d_top + Inches(2.0), d_left + Inches(2.15), d_top + Inches(2.4), color=INST_LIGHT, weight=1.5)

# Legend at bottom
leg_y = d_top + Inches(3.5)
# swatch + text pairs
def leg_row(yoff, fill, border, dashed, text):
    sw = add_card(s, d_left, leg_y + yoff, Inches(0.18), Inches(0.18), fill, border, border_w=1.5)
    if dashed: sw.line.dash_style = 7
    add_text(s, text, d_left + Inches(0.25), leg_y + yoff - Inches(0.04), Inches(6.5), Inches(0.25),
             size=9.5, color=INST_DARK)

leg_row(Inches(0.0), RGBColor(0xE8, 0xF7, 0xEF), INST_GREEN, True,
        "Owned by Deployment Portal (read + write)")
leg_row(Inches(0.28), BLUE_BG, BLUE, False, "Partial integration")
leg_row(Inches(0.56), GREEN_LIGHT_BG, GREEN_LIGHT, False, "No integration yet")

# RIGHT: 3 small workstream cards
r_x = LEFT_M + Inches(7.8)
r_w = SLIDE_W - RIGHT_M - r_x
add_text(s, "Connect next", r_x, TOP_M + Inches(0.95), r_w, Inches(0.3),
         size=14, bold=True, color=INST_DARK)
ws = [
    ("Cert & learning", "Michael", "[BLANK — platform + endpoints]"),
    ("Customer tracking", "Amy", "[BLANK — what + where]"),
    ("Amplitude metrics", "[BLANK — owner]",
     "Track: opens, syncs, writeback success, AI queries, activation"),
]
ws_y = TOP_M + Inches(1.4)
ws_h = Inches(1.5)
ws_gap = Inches(0.15)
for i, (title, owner, body) in enumerate(ws):
    y = ws_y + i*(ws_h + ws_gap)
    sh = add_card(s, r_x, y, r_w, ws_h, WHITE, INST_BORDER, border_w=1.5)
    # accent stripe
    stripe = slide_stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, r_x, y, Inches(0.06), ws_h)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = [INST_GREEN, BLUE, PURPLE][i]
    stripe.line.fill.background()
    _no_shadow(stripe)
    # title
    add_text(s, title, r_x + Inches(0.18), y + Inches(0.12), r_w - Inches(0.3), Inches(0.3),
             size=12, bold=True, color=INST_DARK)
    # owner
    add_text(s, owner, r_x + Inches(0.18), y + Inches(0.42), r_w - Inches(0.3), Inches(0.25),
             size=9.5, italic=True, color=INST_MUTED)
    # body
    body_tb = s.shapes.add_textbox(r_x + Inches(0.18), y + Inches(0.7),
                                    r_w - Inches(0.3), ws_h - Inches(0.75))
    body_tf = body_tb.text_frame
    body_tf.margin_left = body_tf.margin_right = body_tf.margin_top = body_tf.margin_bottom = 0
    body_tf.word_wrap = True
    body_p = body_tf.paragraphs[0]
    _render_text_with_blanks(body_p, body, size=10)

# Save
out = '/Users/asang.mehta/Downloads/firebase-project/team-presentation-v4.5.3.pptx'
prs.save(out)

# Post-process: coerce decimal EMU coordinates to integers.
# python-pptx serializes float arithmetic results as "123456.0" which is valid
# OOXML-ish (PowerPoint forgives it), but Google Slides' stricter parser
# rejects the file with an "unable to open" error.
# Affected attributes: x, y, cx, cy, off (anything that's an EMU long).
import zipfile, shutil, re, os
def _fix_decimals(path):
    tmp = path + '.tmp'
    pattern = re.compile(r'((?:x|y|cx|cy)=")(\d+)\.\d+(")')
    with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml'):
                text = data.decode('utf-8')
                fixed = pattern.sub(r'\1\2\3', text)
                if fixed != text:
                    data = fixed.encode('utf-8')
            zout.writestr(item, data)
    shutil.move(tmp, path)
_fix_decimals(out)

print(f"✓ Saved: {out}")
print(f"  Size: {os.path.getsize(out)/1024:.1f} KB")
print(f"  Slides: {len(prs.slides)}")

# Verify no decimals remain
import subprocess
result = subprocess.run(['unzip', '-p', out, 'ppt/slides/slide1.xml'], capture_output=True, text=True)
remaining = re.findall(r'(?:x|y|cx|cy)="\d+\.\d+"', result.stdout)
print(f"  Decimal coords remaining: {len(remaining)} (should be 0)")
